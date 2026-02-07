"""
@deprecated このモジュールは非推奨です。

理由:
    システムは画像ベースのベクトル検索を採用しており、
    テキストチャンク抽出機能は使用されていません。
    
    実際の処理フロー:
    ドキュメント → 画像変換 → OCI Embedding API → DB保存（IMG_EMBEDDINGS）
    
    画像ベースの処理は parallel_processor.py と image_vectorizer.py で実装されています。

注意:
    将来的にテキストベースの検索を実装する予定はありません。
    このファイルは後方互換性のために残されていますが、新規開発では使用しないでください。
"""
import asyncio
import base64
import logging
import subprocess
import tempfile
import uuid
import warnings
from pathlib import Path
from typing import Any, Dict

import PyPDF2
from PIL import Image
from pptx import Presentation

logger = logging.getLogger(__name__)

# 非推奨警告を発行
warnings.warn(
    "document_processor モジュールは非推奨です。"
    "画像ベースの処理には parallel_processor.py と image_vectorizer.py を使用してください。",
    DeprecationWarning,
    stacklevel=2
)


class DocumentProcessor:
    """文書処理サービス - PDF/DOCX/DOC/PPT/PPTX/PNG/JPG/JPEGの解析"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_doc,
            'pptx': self._process_pptx,
            'ppt': self._process_ppt,
            'png': self._process_image,
            'jpg': self._process_image,
            'jpeg': self._process_image,
        }
    
    def _process_pdf(self, file_path: str, filename: str) -> Dict[str, Any]:
        """PDF処理"""
        chunks = []
        
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            page_count = len(reader.pages)
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    # ページごとにチャンク化
                    chunks.append({
                        "page_number": page_num,
                        "text": text.strip(),
                        "chunk_id": f"{uuid.uuid4()}"
                    })
        
        return {
            "chunks": chunks,
            "page_count": page_count,
            "metadata": {
                "format": "PDF",
                "total_pages": page_count
            }
        }
    
    def _process_pptx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """PowerPoint処理"""
        chunks = []
        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts = []
            
            # スライド内の全テキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            
            slide_text = "\n".join(text_parts).strip()
            
            if slide_text:
                chunks.append({
                    "page_number": slide_num,
                    "text": slide_text,
                    "chunk_id": f"{uuid.uuid4()}"
                })
        
        return {
            "chunks": chunks,
            "page_count": slide_count,
            "metadata": {
                "format": "PowerPoint",
                "total_slides": slide_count
            }
        }
    
    def _convert_with_libreoffice(self, file_path: str, filename: str, doc_format: str, doc_type: str = "Word文書") -> Dict[str, Any]:
        """
        LibreOfficeを使用してOffice文書を処理する共通メソッド
        
        Args:
            file_path: 入力ファイルパス
            filename: ファイル名
            doc_format: ドキュメント形式（'.docx', '.doc', '.ppt'など）
            doc_type: ドキュメントタイプ（"Word文書", "プレゼンテーション"など）
        
        Returns:
            処理結果（チャンク、ページ数、メタデータ）
        """
        try:
            # 一時ディレクトリを作成
            with tempfile.TemporaryDirectory() as temp_dir:
                # LibreOfficeでPDFに変換
                result = subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    file_path
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode != 0:
                    raise ValueError(f"LibreOffice変換エラー: {result.stderr}")
                
                # 変換されたPDFファイルを検索（元のファイル名をベースに生成される）
                original_name = Path(file_path).stem
                pdf_path = Path(temp_dir) / f"{original_name}.pdf"
                
                if not pdf_path.exists():
                    # フォールバック: ディレクトリ内のPDFファイルを検索
                    pdf_files = list(Path(temp_dir).glob("*.pdf"))
                    if pdf_files:
                        pdf_path = pdf_files[0]
                        logger.info(f"フォールバックPDFファイル使用: {pdf_path.name}")
                    else:
                        raise ValueError("PDF変換ファイルが生成されませんでした")
                
                logger.info(f"LibreOffice変換成功: {filename} -> {pdf_path.name}")
                
                # PyPDF2でPDFからテキスト抽出
                chunks = []
                with open(pdf_path, 'rb') as pdf_file:
                    reader = PyPDF2.PdfReader(pdf_file)
                    total_pages = len(reader.pages)
                    
                    for page_num, page in enumerate(reader.pages, start=1):
                        text = page.extract_text()
                        if text.strip():
                            # ページごとにチャンク化
                            chunks.append({
                                "page_number": page_num,
                                "text": text.strip(),
                                "chunk_id": f"{uuid.uuid4()}"
                            })
                
                if not chunks:
                    # テキストがない場合
                    chunks = [{
                        "page_number": 1,
                        "text": f"[{doc_type}: {filename}]\nテキスト内容が見つかりません。",
                        "chunk_id": f"{uuid.uuid4()}"
                    }]
                    total_pages = 1
                
                return {
                    "chunks": chunks,
                    "page_count": total_pages,
                    "metadata": {
                        "format": f"{doc_type} ({doc_format})",
                        "extraction_method": "LibreOffice + PyPDF2",
                        "conversion_quality": "high",
                        "total_pages": total_pages,
                        "extracted_chunks": len(chunks)
                    }
                }
            
        except subprocess.TimeoutExpired:
            logger.error(f"LibreOfficeタイムアウト: {filename}")
            return {
                "chunks": [{
                    "page_number": 1,
                    "text": f"[{doc_type}: {filename}]\n処理がタイムアウトしました。",
                    "chunk_id": f"{uuid.uuid4()}"
                }],
                "page_count": 1,
                "metadata": {
                    "format": f"{doc_type} ({doc_format})",
                    "error": "timeout"
                }
            }
        except Exception as e:
            logger.error(f"{doc_format}ファイル処理エラー: {filename} - {e}")
            return {
                "chunks": [{
                    "page_number": 1,
                    "text": f"[{doc_type}: {filename}]\n処理エラー: {str(e)}",
                    "chunk_id": f"{uuid.uuid4()}"
                }],
                "page_count": 1,
                "metadata": {
                    "format": f"{doc_type} ({doc_format})",
                    "error": str(e)
                }
            }
    
    def _process_docx(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Word文書(.docx)処理 - LibreOfficeを使用"""
        return self._convert_with_libreoffice(file_path, filename, ".docx")
    
    def _process_doc(self, file_path: str, filename: str) -> Dict[str, Any]:
        """Word文書(.doc)処理 - LibreOfficeを使用"""
        return self._convert_with_libreoffice(file_path, filename, ".doc")
    
    def _process_ppt(self, file_path: str, filename: str) -> Dict[str, Any]:
        """
        PowerPoint(.ppt)処理 - LibreOfficeを使用
        
        Note:
            python-pptxは.ppt形式をサポートしていないため、
            LibreOfficeでPDFに変換してからテキスト抽出を行います。
        """
        return self._convert_with_libreoffice(file_path, filename, ".ppt", "プレゼンテーション")
    
    def _process_txt(self, file_path: str, filename: str) -> Dict[str, Any]:
        """
        テキストファイル処理
        
        Note:
            このメソッドは現在supported_formatsに登録されていません。
            将来的にtxt/md形式のサポートを追加する際に使用します。
        """
        chunks = []
        
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # 適切なサイズにチャンク化（1000文字程度）
        chunk_size = 1000
        text_chunks = []
        
        for i in range(0, len(content), chunk_size):
            chunk_text = content[i:i + chunk_size].strip()
            if chunk_text:
                text_chunks.append(chunk_text)
        
        for idx, chunk_text in enumerate(text_chunks, start=1):
            chunks.append({
                "page_number": idx,
                "text": chunk_text,
                "chunk_id": f"{uuid.uuid4()}"
            })
        
        return {
            "chunks": chunks,
            "page_count": len(chunks),
            "metadata": {
                "format": "Text",
                "total_chars": len(content)
            }
        }
    
    def _process_image(self, file_path: str, filename: str) -> Dict[str, Any]:
        """画像ファイル処理 - OCI Vision AIを使用してテキスト抽出"""
        try:
            from app.services.ai_copilot import get_copilot_service
            
            # 画像の検証とサイズ取得（with文でリソース管理）
            try:
                with Image.open(file_path) as img:
                    width, height = img.size
                    logger.info(f"画像サイズ: {width}x{height}")
            except Exception as e:
                logger.error(f"画像の検証エラー: {e}")
                raise ValueError(f"無効な画像ファイル: {e}")
            
            # 画像をbase64エンコード
            with open(file_path, 'rb') as f:
                image_data = f.read()
            base64_image = base64.b64encode(image_data).decode('utf-8')
            
            # MIMEタイプを判定
            file_ext = Path(filename).suffix.lower().lstrip('.')
            mime_type = f"image/{file_ext}" if file_ext in ['png', 'jpg', 'jpeg'] else 'image/jpeg'
            data_url = f"data:{mime_type};base64,{base64_image}"
            
            # Vision AIでテキスト抽出
            copilot = get_copilot_service()
            prompt = "この画像に含まれるすべてのテキストを抽出してください。テキストがない場合は「テキストなし」と応答してください。"
            
            # 非同期テキスト抽出
            async def extract_text():
                result = []
                async for chunk in copilot.chat_stream(
                    message=prompt,
                    images=[{"data_url": data_url}]
                ):
                    result.append(chunk)
                return ''.join(result)
            
            # asyncio.run()を使用（Python 3.7+推奨方式）
            try:
                extracted_text = asyncio.run(extract_text())
            except RuntimeError:
                # 既存のイベントループ内で実行されている場合のフォールバック
                loop = asyncio.get_event_loop()
                extracted_text = loop.run_until_complete(extract_text())
            
            logger.info(f"画像からテキストを抽出: {len(extracted_text)}文字")
            
            # テキストがない場合
            if not extracted_text or extracted_text.strip().lower() in ['テキストなし', 'no text']:
                extracted_text = f"[画像: {filename}]\n画像からテキストを抽出できませんでした。"
            
            # チャンクとして返す
            chunks = [{
                "page_number": 1,
                "text": extracted_text.strip(),
                "chunk_id": f"{uuid.uuid4()}"
            }]
            
            return {
                "chunks": chunks,
                "page_count": 1,
                "metadata": {
                    "format": "Image",
                    "image_size": f"{width}x{height}",
                    "mime_type": mime_type
                }
            }
            
        except Exception as e:
            logger.error(f"画像処理エラー: {filename} - {e}")
            # エラーの場合でも基本情報を返す
            return {
                "chunks": [{
                    "page_number": 1,
                    "text": f"[画像: {filename}]\n画像処理エラー: {str(e)}",
                    "chunk_id": f"{uuid.uuid4()}"
                }],
                "page_count": 1,
                "metadata": {
                    "format": "Image",
                    "error": str(e)
                }
            }

# シングルトンインスタンス
document_processor = DocumentProcessor()
