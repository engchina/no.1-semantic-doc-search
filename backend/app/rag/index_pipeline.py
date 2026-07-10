from __future__ import annotations

import asyncio
import hashlib
import json
import logging
import mimetypes
import os
import re
import tempfile
import unicodedata
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path, PurePosixPath
from typing import Any
from uuid import uuid4
from xml.etree import ElementTree

from docx import Document
from PIL import Image
from PyPDF2 import PdfReader
from pptx import Presentation

from app.rag.clients import embedding_client, mineru_client, ocr_client, vlm_client
from app.rag.models import ProfileConfig, VlmExtractionOutput
from app.rag.oracle_repository import EvidenceRecord, VlmFacetRecord, rag_repository
from app.rag.profile_repository import profile_repository
from app.rag.service_settings import retrieval_service_settings
from app.services.oci_service import oci_service
from app.services.parallel_processor import _convert_file_to_images_worker

logger = logging.getLogger(__name__)

INDEX_OUTPUT_CONTRACT = """Return one JSON object with exactly this technical shape:
{
  "summary": "concise source-grounded summary",
  "keywords": ["searchable term"],
  "facts": [{"text": "source-grounded fact", "source_locator": "page:N", "confidence": 0.0}]
}
Omit unsupported or uncertain facts. Do not add any other keys."""


@dataclass
class SourceBlock:
    page_number: int
    text: str
    kind: str
    source: str
    bbox: list[float] | None = None


@dataclass
class PageExtraction:
    page_number: int
    image: bytes | None = None
    image_path: Path | None = None
    image_dpi: int | None = None
    asset_object_name: str | None = None
    native_text: str = ""
    mineru_blocks: list[SourceBlock] = field(default_factory=list)
    ocr_blocks: list[SourceBlock] = field(default_factory=list)
    ocr_engine: str | None = None


@dataclass
class IndexOutcome:
    object_name: str
    document_id: str
    matched_profiles: list[int]
    indexed_profiles: list[int]
    reused_profiles: list[int]
    failed_profiles: dict[int, str]
    page_count: int
    degraded_services: list[str]


def _clean_text(value: object) -> str:
    text = unicodedata.normalize("NFC", str(value or "").replace("\x00", ""))
    cjk_gap = re.compile(r"(?<=[ぁ-んァ-ン一-龯々])[ \t　]+(?=[ぁ-んァ-ン一-龯々])")
    previous = None
    while previous != text:
        previous = text
        text = cjk_gap.sub("", text)
    text = re.sub(r"(?<=[ァ-ン])一(?=[ァ-ン])", "ー", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


def _page_image(page: PageExtraction) -> bytes | None:
    if page.image is not None:
        return page.image
    if page.image_path and page.image_path.exists():
        return page.image_path.read_bytes()
    return None


def _image_at_dpi(image: bytes, source_dpi: int | None, target_dpi: int) -> bytes:
    if source_dpi is None or target_dpi >= source_dpi:
        return image
    with Image.open(BytesIO(image)) as source:
        size = tuple(max(1, round(value * target_dpi / source_dpi)) for value in source.size)
        resized = source.resize(size, Image.Resampling.LANCZOS)
        output = BytesIO()
        resized.save(output, format="PNG")
        return output.getvalue()


def _xlsx_pages(content: bytes) -> dict[int, str]:
    namespace = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    with zipfile.ZipFile(BytesIO(content)) as archive:
        shared: list[str] = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
            shared = [
                "".join(node.text or "" for node in item.findall(".//m:t", namespace))
                for item in root.findall("m:si", namespace)
            ]
        sheets = sorted(
            name for name in archive.namelist()
            if re.fullmatch(r"xl/worksheets/sheet\d+\.xml", name)
        )
        pages: dict[int, str] = {}
        for page_number, name in enumerate(sheets, start=1):
            root = ElementTree.fromstring(archive.read(name))
            rows: list[str] = []
            for row in root.findall(".//m:row", namespace):
                values: list[str] = []
                for cell in row.findall("m:c", namespace):
                    node = cell.find("m:v", namespace)
                    value = node.text if node is not None else ""
                    if cell.attrib.get("t") == "s" and value and value.isdigit():
                        index = int(value)
                        value = shared[index] if index < len(shared) else value
                    values.append(value or "")
                rows.append("\t".join(values))
            pages[page_number] = _clean_text("\n".join(rows))
        return pages


def _native_pages(content: bytes, extension: str) -> dict[int, str]:
    try:
        if extension == "pdf":
            return {
                index: _clean_text(page.extract_text() or "")
                for index, page in enumerate(PdfReader(BytesIO(content)).pages, start=1)
            }
        if extension == "pptx":
            pages: dict[int, str] = {}
            for page_number, slide in enumerate(Presentation(BytesIO(content)).slides, start=1):
                values: list[str] = []
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        values.append(shape.text)
                    if getattr(shape, "has_table", False):
                        values.extend("\t".join(cell.text for cell in row.cells) for row in shape.table.rows)
                pages[page_number] = _clean_text("\n".join(values))
            return pages
        if extension == "docx":
            document = Document(BytesIO(content))
            values = [paragraph.text for paragraph in document.paragraphs]
            values.extend(
                "\t".join(cell.text for cell in row.cells)
                for table in document.tables for row in table.rows
            )
            return {1: _clean_text("\n".join(values))}
        if extension == "xlsx":
            return _xlsx_pages(content)
        if extension in {"txt", "md", "csv", "tsv", "html", "htm"}:
            return {1: _clean_text(content.decode("utf-8", errors="replace"))}
    except Exception as error:
        logger.warning("Native text extraction failed for %s: %s", extension, error)
    return {}


def _bbox(value: object) -> list[float] | None:
    if not isinstance(value, list) or len(value) != 4:
        return None
    try:
        result = [float(item) for item in value]
    except (TypeError, ValueError):
        return None
    return result if result[2] >= result[0] and result[3] >= result[1] else None


def _block_text(value: dict[str, Any]) -> str:
    parts: list[str] = []
    for key in (
        "text", "table_body", "content", "markdown", "latex", "equation",
        "caption", "image_caption", "image_footnote",
    ):
        candidate = value.get(key)
        if isinstance(candidate, list):
            candidate = " ".join(str(item) for item in candidate if item)
        if isinstance(candidate, str) and candidate.strip():
            parts.append(candidate.strip())
    return _clean_text("\n".join(dict.fromkeys(parts)))


def _mineru_blocks(result: dict[str, Any]) -> list[SourceBlock]:
    blocks: list[SourceBlock] = []
    for item in mineru_client.content_blocks(result):
        raw_page = item.get("page_idx", item.get("page_number", 0))
        try:
            page_number = int(raw_page) + (1 if "page_idx" in item else 0)
        except (TypeError, ValueError):
            page_number = 1
        text = _block_text(item)
        if text or str(item.get("type") or "").casefold() == "image":
            blocks.append(
                SourceBlock(
                    page_number=max(1, page_number),
                    text=text,
                    kind=str(item.get("type") or item.get("category") or "block")[:40],
                    source="mineru",
                    bbox=_bbox(item.get("bbox")),
                )
            )
    return blocks


def _mineru_missing_pages(pages: list[PageExtraction]) -> list[PageExtraction]:
    return [page for page in pages if not any(block.text for block in page.mineru_blocks)]


def _ocr_blocks(result: dict[str, Any], page_number: int) -> list[SourceBlock]:
    blocks = [
        SourceBlock(
            page_number=page_number,
            text=_clean_text(cell.get("text")),
            kind=str(cell.get("category") or "ocr_block")[:40],
            source=str(result.get("engine") or "ocr"),
            bbox=_bbox(cell.get("bbox")),
        )
        for cell in (result.get("cells") or [])
        if isinstance(cell, dict) and _clean_text(cell.get("text"))
    ]
    text = _clean_text(result.get("text"))
    if not blocks and text:
        blocks.append(SourceBlock(page_number, text, "ocr_text", str(result.get("engine") or "ocr")))
    return blocks


def _chunks(text: str, target: int = 1800, overlap: int = 200) -> list[str]:
    text = _clean_text(text)
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(len(text), start + target)
        if end < len(text):
            boundary = max(text.rfind("\n", start + target // 2, end), text.rfind("。", start, end))
            if boundary > start:
                end = boundary + 1
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return [value for value in chunks if value]


async def _run_ocr(
    page: PageExtraction,
    degraded: list[str],
    semaphores: dict[str, asyncio.Semaphore] | None = None,
) -> None:
    image = _page_image(page)
    if image is None:
        return
    settings = retrieval_service_settings.get_ocr(mask_secrets=False)
    engines = (("dots", settings.dots), ("glm", settings.glm), ("unlimited", settings.unlimited))
    semaphores = semaphores or {name: asyncio.Semaphore(value.workers) for name, value in engines}
    dots_partial: list[SourceBlock] = []
    for name, engine in engines:
        if not engine.enabled:
            continue
        try:
            async with semaphores[name]:
                result = await ocr_client.recognize(
                    engine=name,
                    settings=engine,
                    image=_image_at_dpi(image, page.image_dpi, engine.dpi),
                )
            blocks = _ocr_blocks(result, page.page_number)
            if not blocks:
                continue
            categories = {
                str(cell.get("category") or "")
                for cell in (result.get("cells") or []) if isinstance(cell, dict)
            }
            if name == "dots" and "Picture" in categories and not categories.intersection(
                {"Text", "Table", "List-item", "Formula"}
            ):
                dots_partial = blocks
                degraded.append("ocr:dots_partial")
                continue
            page.ocr_blocks, page.ocr_engine = blocks, name
            return
        except Exception as error:
            degraded.append(f"ocr:{name}")
            logger.warning("%s OCR failed on page %s: %s", name, page.page_number, error)
    if dots_partial:
        page.ocr_blocks, page.ocr_engine = dots_partial, "dots_partial"


def _shared_page_text(page: PageExtraction) -> tuple[str, list[SourceBlock]]:
    blocks = [*page.mineru_blocks, *page.ocr_blocks]
    values = [page.native_text, *(block.text for block in blocks)]
    unique = [value for value in dict.fromkeys(_clean_text(item) for item in values) if value]
    return "\n\n".join(unique), blocks


async def _build_shared_evidence(*, file_name: str, pages: list[PageExtraction]) -> list[EvidenceRecord]:
    evidence: list[EvidenceRecord] = []
    text_cache: dict[str, list[float]] = {}
    visual_cache: dict[int, list[float]] = {}
    for page in pages:
        page_text, blocks = _shared_page_text(page)
        locator = f"page:{page.page_number}"
        parent = EvidenceRecord(
            evidence_id=uuid4().hex,
            document_id="",
            page_number=page.page_number,
            unit_kind="page",
            source_locator=locator,
            raw_text=page_text,
            search_text=_clean_text(f"{file_name}\n{page_text}") or file_name,
            asset_object_name=page.asset_object_name,
            provenance={"native": bool(page.native_text), "ocr_engine": page.ocr_engine},
        )
        evidence.append(parent)
        for index, block in enumerate(blocks, start=1):
            evidence.append(
                EvidenceRecord(
                    evidence_id=uuid4().hex,
                    document_id="",
                    parent_evidence_id=parent.evidence_id,
                    page_number=page.page_number,
                    unit_kind=block.kind[:40] or "block",
                    source_locator=f"{locator}/block:{index}",
                    bbox=block.bbox,
                    raw_text=block.text,
                    search_text=_clean_text(f"{file_name}\n{block.text}") or file_name,
                    asset_object_name=page.asset_object_name,
                    provenance={"source": block.source},
                )
            )
        for index, chunk in enumerate(_chunks(page_text), start=1):
            evidence.append(
                EvidenceRecord(
                    evidence_id=uuid4().hex,
                    document_id="",
                    parent_evidence_id=parent.evidence_id,
                    page_number=page.page_number,
                    unit_kind="text_chunk",
                    source_locator=f"{locator}/chunk:{index}",
                    raw_text=chunk,
                    search_text=_clean_text(f"{file_name}\n{chunk}") or file_name,
                    asset_object_name=page.asset_object_name,
                    provenance={"source": "merged_text"},
                )
            )
    for item in evidence:
        cache_key = hashlib.sha256(item.search_text[:6000].encode()).hexdigest()
        if cache_key not in text_cache:
            text_cache[cache_key] = (await embedding_client.text([item.search_text[:6000]]))[0]
        item.text_embedding = text_cache[cache_key]
        if item.unit_kind == "page" and item.asset_object_name:
            page = next(value for value in pages if value.page_number == item.page_number)
            image = _page_image(page)
            if image is not None:
                if page.page_number not in visual_cache:
                    visual_cache[page.page_number] = await embedding_client.image(image, "image/png")
                item.visual_embedding = visual_cache[page.page_number]
    return evidence


def _page_evidence(evidence: list[EvidenceRecord]) -> list[EvidenceRecord]:
    return [item for item in evidence if item.unit_kind == "page" and item.page_number is not None]


async def _build_profile_facets(
    *, profile: ProfileConfig, object_name: str,
    evidence: list[EvidenceRecord], page_images: dict[int, bytes],
    on_page=None,
) -> list[VlmFacetRecord]:
    facets: list[VlmFacetRecord] = []
    for item in _page_evidence(evidence):
        prompt = (
            f"Administrator extraction instruction:\n{profile.extraction_prompt}\n\n"
            f"Document context: {json.dumps({'object_name': object_name, 'page_number': item.page_number}, ensure_ascii=False)}\n"
            f"Page text:\n{item.raw_text[:12000]}\n\nSource locator: {item.source_locator}\n\n"
            f"{INDEX_OUTPUT_CONTRACT}"
        )
        output = VlmExtractionOutput.model_validate(
            await vlm_client.generate_json(prompt=prompt, image=page_images.get(item.page_number))
        )
        if output.search_text():
            facets.append(
                VlmFacetRecord(
                    evidence_id=item.evidence_id,
                    output=output,
                    text_embedding=(await embedding_client.text([output.search_text()[:6000]]))[0],
                )
            )
        if on_page:
            on_page()
    return facets


class SharedIndexPipeline:
    async def index_object(
        self, object_name: str, *, force: bool = False, profile_slots: set[int] | None = None,
        progress=None,
    ) -> IndexOutcome:
        """progress: callable(done_units, total_units)。VLM抽出1ページ完了ごとに呼ばれる。"""
        content = await asyncio.to_thread(oci_service.download_object, object_name)
        if not content:
            raise FileNotFoundError(f"Object Storage object is unavailable: {object_name}")
        path = PurePosixPath(object_name)
        extension = path.suffix.casefold().lstrip(".")
        file_name = path.name
        media_type = mimetypes.guess_type(file_name)[0] or "application/octet-stream"
        document = rag_repository.upsert_document(
            bucket=os.environ.get("OCI_BUCKET") or "",
            object_name=object_name,
            file_name=file_name,
            media_type=media_type,
            content=content,
            document_type=extension or None,
        )
        profiles = [
            profile for profile in profile_repository.enabled_profiles()
            if profile_slots is None or profile.slot_no in profile_slots
        ]
        vlm_model = oci_service.get_enterprise_ai_settings().model or ""
        mineru = retrieval_service_settings.get_mineru()
        ocr = retrieval_service_settings.get_ocr(mask_secrets=False)
        embedding_model = os.environ.get("OCI_COHERE_EMBED_MODEL", "cohere.embed-v4.0")
        config_hash = hashlib.sha256(
            json.dumps(
                {
                    "pipeline": "shared-index-v2",
                    "native_parser": "native+page-render-v3",
                    "mineru": mineru.model_dump(mode="json"),
                    "ocr": {
                        "enabled": ocr.enabled,
                        "engines": [
                            (name, value.enabled, value.base_url, value.model, value.dpi)
                            for name, value in (("dots", ocr.dots), ("glm", ocr.glm), ("unlimited", ocr.unlimited))
                        ],
                    },
                    "embedding_model": embedding_model,
                },
                ensure_ascii=False,
                sort_keys=True,
            ).encode()
        ).hexdigest()
        index_run_id = None if force else rag_repository.reusable_document_run(
            document_id=document.document_id,
            content_sha256=document.content_sha256,
            config_hash=config_hash,
        )
        degraded: list[str] = []
        uploaded_assets: list[str] = []
        if index_run_id:
            index_run_id, evidence = rag_repository.serving_evidence(document.document_id)
            page_images: dict[int, bytes] = {}
            for item in evidence:
                if item.unit_kind == "page" and item.page_number and item.asset_object_name:
                    image = await asyncio.to_thread(oci_service.download_object, item.asset_object_name)
                    if image:
                        page_images[item.page_number] = image
            page_count = sum(item.unit_kind == "page" for item in evidence)
            rag_repository.set_document_status(document.document_id, "INDEXED")
        else:
            index_run_id = rag_repository.start_index_run(
                document_id=document.document_id,
                content_sha256=document.content_sha256,
                config_hash=config_hash,
                native_parser="native+page-render-v3",
                embedding_model=embedding_model,
            )
            render_temp = tempfile.TemporaryDirectory(prefix="rag-render-")
            try:
                native = await asyncio.to_thread(_native_pages, content, extension)
                enabled_engines = [value for value in (ocr.dots, ocr.glm, ocr.unlimited) if value.enabled]
                render_dpi = max([200, *(value.dpi for value in enabled_engines)]) if ocr.enabled else 200
                rendered_paths: dict[int, Path] = {}
                ranges: list[tuple[int | None, int | None]] = [(None, None)]
                if extension == "pdf" and native:
                    ranges = [(start, min(len(native), start + 19)) for start in range(1, len(native) + 1, 20)]
                for first_page, last_page in ranges:
                    success, rendered, error = await asyncio.to_thread(
                        _convert_file_to_images_worker, content, extension, object_name,
                        render_dpi, first_page, last_page,
                    )
                    if not success:
                        degraded.append("page_render")
                        logger.warning("Page rendering failed for %s: %s", object_name, error)
                        continue
                    for page_number, image in rendered:
                        local_path = Path(render_temp.name) / f"page_{page_number:06d}.png"
                        local_path.write_bytes(image)
                        rendered_paths[page_number] = local_path
                page_numbers = sorted(set(native) | set(rendered_paths)) or [1]
                pages = [
                    PageExtraction(
                        page_number=number,
                        image_path=rendered_paths.get(number),
                        image_dpi=None if extension in {"png", "jpg", "jpeg"} else render_dpi,
                        native_text=native.get(number, ""),
                    )
                    for number in page_numbers
                ]
                folder = str(path.with_suffix(""))
                for page in pages:
                    image = _page_image(page)
                    if image is None:
                        continue
                    page.asset_object_name = f"{folder}/page_{page.page_number:03d}_{index_run_id}.png"
                    uploaded = await asyncio.to_thread(
                        oci_service.upload_file, image, page.asset_object_name, "image/png",
                        f"page_{page.page_number:03d}.png", len(image),
                    )
                    if uploaded:
                        uploaded_assets.append(page.asset_object_name)
                    else:
                        page.asset_object_name = None
                        degraded.append("asset_upload")
                mineru_version: str | None = None
                if mineru.enabled and mineru.base_url:
                    try:
                        result = await mineru_client.parse_file(
                            file_name=file_name, content=content, media_type=media_type, settings=mineru
                        )
                        mineru_version = str(result.get("version") or result.get("engine_version") or "") or None
                        for block in _mineru_blocks(result):
                            page = next((value for value in pages if value.page_number == block.page_number), None)
                            if page:
                                page.mineru_blocks.append(block)
                    except Exception as error:
                        degraded.append("mineru")
                        logger.warning("MinerU failed for %s: %s", object_name, error)
                if ocr.enabled:
                    missing_mineru = _mineru_missing_pages(pages)
                    semaphores = {
                        name: asyncio.Semaphore(value.workers)
                        for name, value in (("dots", ocr.dots), ("glm", ocr.glm), ("unlimited", ocr.unlimited))
                    }
                    await asyncio.gather(*(_run_ocr(page, degraded, semaphores) for page in missing_mineru))
                evidence = await _build_shared_evidence(file_name=file_name, pages=pages)
                for item in evidence:
                    item.document_id = document.document_id
                used_ocr = sorted({page.ocr_engine for page in pages if page.ocr_engine})
                covered = sum(bool(page.native_text or page.mineru_blocks or page.ocr_blocks) for page in pages)
                rag_repository.store_document_evidence(
                    index_run_id=index_run_id,
                    document_id=document.document_id,
                    evidence=evidence,
                    page_count=len(pages),
                    page_coverage=covered / max(1, len(pages)),
                    mineru_version=mineru_version,
                    ocr_engines=used_ocr,
                )
                page_images = {
                    page.page_number: image for page in pages if (image := _page_image(page)) is not None
                }
                page_count = len(pages)
            except Exception as error:
                rag_repository.fail_index_run(index_run_id, document.document_id, str(error))
                for asset in uploaded_assets:
                    await asyncio.to_thread(oci_service.delete_object, asset)
                raise
            finally:
                render_temp.cleanup()

        indexed: list[int] = []
        reused: list[int] = []
        failed: dict[int, str] = {}
        page_units = len(_page_evidence(evidence))
        total_units = max(1, page_units * len(profiles))
        done_units = 0

        def _bump(count: int = 1) -> None:
            nonlocal done_units
            done_units += count
            if progress:
                progress(min(done_units, total_units), total_units)

        if progress:
            progress(0, total_units)
        for profile in profiles:
            profile_runtime_hash = hashlib.sha256(
                f"{profile.config_hash or ''}:{vlm_model}".encode()
            ).hexdigest()
            if rag_repository.reusable_profile_run(
                document_id=document.document_id,
                profile=profile,
                index_run_id=index_run_id,
                content_sha256=document.content_sha256,
                config_hash=profile_runtime_hash,
            ):
                reused.append(profile.slot_no)
                _bump(page_units)
                continue
            profile_repository.set_apply_status(profile.slot_no, "PROCESSING")
            try:
                facets = await _build_profile_facets(
                    profile=profile, object_name=object_name, evidence=evidence, page_images=page_images,
                    on_page=_bump,
                )
                rag_repository.store_profile_facets(
                    document_id=document.document_id,
                    index_run_id=index_run_id,
                    content_sha256=document.content_sha256,
                    profile=profile,
                    facets=facets,
                    config_hash=profile_runtime_hash,
                )
                indexed.append(profile.slot_no)
            except Exception as error:
                message = str(error)[:2000]
                failed[profile.slot_no] = message
                rag_repository.record_profile_failure(
                    document_id=document.document_id,
                    index_run_id=index_run_id,
                    content_sha256=document.content_sha256,
                    profile=profile,
                    error=message,
                    config_hash=profile_runtime_hash,
                )
                logger.exception("VLM Profile %s indexing failed", profile.slot_no)
        for profile in profiles:
            if profile.slot_no in failed:
                profile_repository.set_apply_status(profile.slot_no, "FAILED")
            else:
                profile_repository.refresh_apply_status(profile.slot_no)
        return IndexOutcome(
            object_name=object_name,
            document_id=document.document_id,
            matched_profiles=[profile.slot_no for profile in profiles],
            indexed_profiles=indexed,
            reused_profiles=reused,
            failed_profiles=failed,
            page_count=page_count,
            degraded_services=sorted(set(degraded)),
        )


index_pipeline = SharedIndexPipeline()
